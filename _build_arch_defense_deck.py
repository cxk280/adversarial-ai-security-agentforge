"""Generates ARCHITECTURE_DEFENSE.pptx from in-line content.

Run once; output is the deck. Kept in-repo so the deck is reproducible
from text and easy to tweak before defense."""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Style
# ---------------------------------------------------------------------------

NAVY = RGBColor(0x1A, 0x2B, 0x4C)
ACCENT = RGBColor(0xC5, 0x3A, 0x3A)
DIM = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(12), Inches(1.5))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AgentForge — Adversarial AI Security Platform"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sub = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(12), Inches(0.7))
    p = sub.text_frame.paragraphs[0]
    p.text = "Multi-agent adversarial evaluation of the OpenEMR Clinical Co-Pilot"
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE

    meta = slide.shapes.add_textbox(Inches(0.7), Inches(6.4), Inches(12), Inches(0.5))
    p = meta.text_frame.paragraphs[0]
    p.text = "Architecture Defense · Gauntlet AI · Week 3 · 2026-05-11 · Chris King"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str | tuple[str, list[str]]],
    *,
    footer: str | None = None,
) -> None:
    """A content slide with a title bar and bullets.

    Each `bullets` entry can be a string (top-level bullet) or a tuple
    (top-level bullet, [sub-bullets])."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.6))
    p = tbox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Body
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.6))
    tf = body.text_frame
    tf.word_wrap = True

    first = True
    for item in bullets:
        if isinstance(item, tuple):
            top, subs = item
        else:
            top, subs = item, []
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = f"• {top}"
        p.font.size = Pt(20)
        p.font.color.rgb = NAVY
        p.space_after = Pt(6)
        for s in subs:
            sub_p = tf.add_paragraph()
            sub_p.text = f"    – {s}"
            sub_p.font.size = Pt(16)
            sub_p.font.color.rgb = DIM
            sub_p.space_after = Pt(2)

    # Footer
    if footer:
        f = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4))
        p = f.text_frame.paragraphs[0]
        p.text = footer
        p.font.size = Pt(11)
        p.font.color.rgb = DIM
        p.alignment = PP_ALIGN.LEFT


def add_diagram_slide(prs: Presentation) -> None:
    """Clean left-to-right pipeline.

    Five boxes in a row: Orchestrator → Red Team → Target → Judge → Documentation.
    Numbered arrows 1-4 between them; one long feedback arrow 5 looping back
    from Judge to Orchestrator below the row.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.6))
    p = tbox.text_frame.paragraphs[0]
    p.text = "Multi-Agent Architecture"
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE

    # Subtitle under title bar
    sub_t = slide.shapes.add_textbox(Inches(0.5), Inches(1.10), Inches(12.3), Inches(0.4))
    sp = sub_t.text_frame.paragraphs[0]
    sp.text = "Five roles, three model families. Attacker and Judge are structurally independent."
    sp.font.size = Pt(15); sp.font.color.rgb = DIM

    # Colors per role
    ORCH_COLOR   = RGBColor(0x1F, 0x3A, 0x5F)
    RED_COLOR    = RGBColor(0xC5, 0x3A, 0x3A)
    TARGET_COLOR = RGBColor(0x4A, 0x55, 0x68)
    JUDGE_COLOR  = RGBColor(0x2E, 0x6E, 0x4A)
    DOC_COLOR    = RGBColor(0x6E, 0x3A, 0x8E)

    # Box geometry — five boxes evenly across, with 0.30" gutters between
    BOX_W = 2.36
    BOX_H = 2.00
    BOX_Y = 2.35
    GUTTER = 0.45
    LEFT_MARGIN = 0.45
    box_x = [LEFT_MARGIN + i * (BOX_W + GUTTER) for i in range(5)]

    def agent_box(idx, name, role_subtitle, model, trust, color):
        x = box_x[idx]
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(BOX_Y), Inches(BOX_W), Inches(BOX_H),
        )
        box.fill.solid(); box.fill.fore_color.rgb = color
        box.line.color.rgb = color
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.14); tf.margin_bottom = Inches(0.14)

        p1 = tf.paragraphs[0]
        p1.text = name
        p1.font.size = Pt(18); p1.font.bold = True; p1.font.color.rgb = WHITE
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(2)

        p2 = tf.add_paragraph()
        p2.text = role_subtitle
        p2.font.size = Pt(11); p2.font.italic = True; p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER
        p2.space_after = Pt(8)

        p3 = tf.add_paragraph()
        p3.text = model
        p3.font.size = Pt(12); p3.font.bold = True; p3.font.color.rgb = WHITE
        p3.alignment = PP_ALIGN.CENTER
        p3.space_after = Pt(4)

        p4 = tf.add_paragraph()
        p4.text = f"trust: {trust}"
        p4.font.size = Pt(11); p4.font.color.rgb = WHITE
        p4.alignment = PP_ALIGN.CENTER

    # ─── 5 boxes ───
    agent_box(0, "Orchestrator", "what to test next",
              "Sonnet 4.6", "high", ORCH_COLOR)
    agent_box(1, "Red Team", "generates + mutates",
              "huihui-ai 70B abl.\n+ DeepSeek-R1", "low", RED_COLOR)
    agent_box(2, "Target", "system under test",
              "Co-Pilot\n(Sonnet 4.6)", "n/a", TARGET_COLOR)
    agent_box(3, "Judge", "scores each attack",
              "Haiku 4.5", "med-hi", JUDGE_COLOR)
    agent_box(4, "Documentation", "writes findings",
              "Sonnet 4.6", "gated", DOC_COLOR)

    # ─── Numbered arrows between adjacent boxes ───
    def arrow_between(i_from: int, i_to: int, num: str, label: str):
        x1 = box_x[i_from] + BOX_W + 0.03
        x2 = box_x[i_to] - 0.03
        y = BOX_Y + BOX_H / 2
        # arrow line
        ln = slide.shapes.add_connector(1, Inches(x1), Inches(y), Inches(x2), Inches(y))
        ln.line.color.rgb = NAVY; ln.line.width = Pt(3)
        from pptx.oxml.ns import qn
        spPr = ln.line._get_or_add_ln()
        head = spPr.makeelement(qn("a:tailEnd"),
                                {"type": "triangle", "w": "med", "len": "med"})
        spPr.append(head)

        # label centered between the boxes, above the arrow
        gap_w = x2 - x1
        cx = (x1 + x2) / 2
        # Number badge
        badge_size = 0.28
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx - badge_size / 2), Inches(y - badge_size - 0.50),
            Inches(badge_size), Inches(badge_size),
        )
        badge.fill.solid(); badge.fill.fore_color.rgb = NAVY
        badge.line.fill.background()
        bp = badge.text_frame.paragraphs[0]
        bp.text = num
        bp.font.size = Pt(10); bp.font.bold = True; bp.font.color.rgb = WHITE
        bp.alignment = PP_ALIGN.CENTER
        badge.text_frame.margin_left = Inches(0)
        badge.text_frame.margin_right = Inches(0)
        badge.text_frame.margin_top = Inches(0.02)
        badge.text_frame.margin_bottom = Inches(0)

        # Label text below the arrow
        lab = slide.shapes.add_textbox(
            Inches(cx - 0.95), Inches(y + 0.05),
            Inches(1.9), Inches(0.40),
        )
        lp = lab.text_frame.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(11); lp.font.bold = True; lp.font.color.rgb = NAVY
        lp.alignment = PP_ALIGN.CENTER

    arrow_between(0, 1, "1", "campaign brief")
    arrow_between(1, 2, "2", "attack (HTTP)")
    arrow_between(2, 3, "3", "response")
    arrow_between(3, 4, "4", "confirmed exploit")

    # ─── Feedback arrow 5: Judge → Orchestrator (loops below) ───
    feedback_y = BOX_Y + BOX_H + 0.95
    judge_cx = box_x[3] + BOX_W / 2
    orch_cx = box_x[0] + BOX_W / 2

    # Down stub from Judge
    s1 = slide.shapes.add_connector(1, Inches(judge_cx), Inches(BOX_Y + BOX_H),
                                    Inches(judge_cx), Inches(feedback_y))
    s1.line.color.rgb = NAVY; s1.line.width = Pt(3); s1.line.dash_style = 7  # dash
    # Horizontal Judge → Orchestrator
    s2 = slide.shapes.add_connector(1, Inches(judge_cx), Inches(feedback_y),
                                    Inches(orch_cx), Inches(feedback_y))
    s2.line.color.rgb = NAVY; s2.line.width = Pt(3); s2.line.dash_style = 7
    # Up arrow to Orchestrator
    s3 = slide.shapes.add_connector(1, Inches(orch_cx), Inches(feedback_y),
                                    Inches(orch_cx), Inches(BOX_Y + BOX_H + 0.03))
    s3.line.color.rgb = NAVY; s3.line.width = Pt(3); s3.line.dash_style = 7
    from pptx.oxml.ns import qn
    spPr = s3.line._get_or_add_ln()
    head = spPr.makeelement(qn("a:tailEnd"),
                            {"type": "triangle", "w": "med", "len": "med"})
    spPr.append(head)

    # Step-5 badge on the feedback line
    cx5 = (judge_cx + orch_cx) / 2
    badge5 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx5 - 0.14), Inches(feedback_y - 0.14),
        Inches(0.28), Inches(0.28),
    )
    badge5.fill.solid(); badge5.fill.fore_color.rgb = NAVY
    badge5.line.fill.background()
    bp = badge5.text_frame.paragraphs[0]
    bp.text = "5"
    bp.font.size = Pt(10); bp.font.bold = True; bp.font.color.rgb = WHITE
    bp.alignment = PP_ALIGN.CENTER
    badge5.text_frame.margin_left = Inches(0)
    badge5.text_frame.margin_right = Inches(0)
    badge5.text_frame.margin_top = Inches(0.02)
    badge5.text_frame.margin_bottom = Inches(0)

    fb_lab = slide.shapes.add_textbox(Inches(cx5 - 2.2), Inches(feedback_y + 0.10),
                                      Inches(4.4), Inches(0.4))
    fp = fb_lab.text_frame.paragraphs[0]
    fp.text = "verdicts + coverage feedback → next campaign"
    fp.font.size = Pt(12); fp.font.bold = True; fp.font.color.rgb = NAVY
    fp.alignment = PP_ALIGN.CENTER

    # ─── Bottom substrate bar ───
    sub = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.40), Inches(6.85),
        Inches(12.55), Inches(0.40),
    )
    sub.fill.solid(); sub.fill.fore_color.rgb = RGBColor(0xF1, 0xF1, 0xF4)
    sub.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0); sub.line.width = Pt(0.5)
    st = sub.text_frame
    st.margin_left = Inches(0.12); st.margin_top = Inches(0.06); st.margin_bottom = Inches(0.06)
    p = st.paragraphs[0]
    p.text = (
        "Substrate: Postgres  ·  Langfuse spans tagged by agent_role  ·  "
        "harness/executor.py enforces target-host allowlist on every HTTP call"
    )
    p.font.size = Pt(11); p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER


def add_section_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(12), Inches(1.2))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    s = slide.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(12), Inches(0.8))
    p = s.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE


# ---------------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------------

def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Title
    add_title_slide(prs)

    # 2. The problem
    add_content_slide(
        prs,
        "The problem",
        [
            "The Clinical Co-Pilot has PHI access, tool-call authority, multi-turn state, and ingests third-party documents.",
            "Manual penetration testing produces a snapshot. The system evolves daily. Static test lists go stale within a sprint.",
            "Failure mode is not 'a jailbreak exists' — it is 'a category of exploit is unaddressed'.",
            ("Required answer: continuous, autonomous adversarial evaluation that adapts as attackers adapt.", [
                "Discover · Generate · Mutate · Judge · Document · Validate · Regress · Report",
            ]),
        ],
        footer="Per Week 3 case study — 'a static test suite is not the answer; an autonomous multi-agent red team is.'",
    )

    # 3. Target at a glance
    add_content_slide(
        prs,
        "Target: OpenEMR Clinical Co-Pilot (built Weeks 1–2)",
        [
            "FastAPI agent, Claude Sonnet 4.6, deployed live on Railway (dev / qa / prod).",
            ("8 patient/FHIR/guideline tools — PHI-bearing, parameter-controlled.", [
                "get_patient_summary, get_medications, get_recent_labs, get_vitals, get_visit_history, get_conditions, search_guidelines, get_extracted_facts",
            ]),
            "LangGraph supervisor: intake_extractor + evidence_retriever + final_answer.",
            ("Existing defenses are mostly soft (prompt-level): system prompt forbids prescriptions, instructs treating retrieved content as DATA.", [
                "Hard defenses: rate limiter (per-session), Pydantic field caps, PHI redaction on traces (NOT responses).",
            ]),
            "Endpoints in scope: /chat, /chat/stream, /chat/graph, /search, /extract, /copilot/lab-trend/{pid}, /copilot/extractions/{pid}, /api/patient-fhir-id/{pid}.",
        ],
    )

    # 4. Threat model — top 5 priorities
    add_content_slide(
        prs,
        "Highest-priority attack categories (from THREAT_MODEL.md)",
        [
            ("1. Authorization bypass — sev 9, priority 7.2", ["Direct HTTP to /copilot/* endpoints from a session bound to a different patient. Best handled deterministically, not by LLM."]),
            ("2. Cross-patient data exfiltration — sev 10, priority 7.0", ["Active-patient constraint is enforced by system-prompt instruction only. Cleanest deterministic signal in the suite."]),
            ("3. Indirect prompt injection — sev 9, priority 6.3", ["Via uploaded PDFs (/extract) and guideline corpus retrieval. System prompt's 'treat retrieved content as DATA' is a soft mitigation, not a control."]),
            ("4. Multi-turn / crescendo injection — sev 8, priority 5.6", ["No turn-level safety re-evaluation; session history is not re-sanitized."]),
            ("5. Persona hijack — clinical authority — sev 10, priority 5.0", ["Producing prescription-shaped text is the worst-case clinical impact. Anthropic alignment is the main current defense."]),
        ],
        footer="Full ranked list of 17 subcategories with severity × exploitability scoring in THREAT_MODEL.md.",
    )

    # 5. Architecture diagram
    add_diagram_slide(prs)

    # 6. Why four agents, three model families
    add_content_slide(
        prs,
        "Why four distinct agents, three model families",
        [
            ("Red Team and Judge must be independent.", [
                "A model that generates an attack is the worst possible evaluator of it — it rationalizes success.",
                "Different model family + different repo path + different prompt template enforces independence.",
            ]),
            ("Red Team and Judge need different alignment.", [
                "Red Team needs an UNCENSORED model — Claude/GPT/Gemini refuse offensive generation at scale.",
                "Judge benefits from FRONTIER alignment — judgment, not generation, so refusals don't apply.",
            ]),
            ("Orchestrator is strategic, not generative — no refusals; Sonnet handles coverage reasoning well.", []),
            ("Documentation is structured prose — Sonnet writes engineer-ready vuln reports.", []),
            "A single-agent or linear-pipeline system fails the spec's explicit requirement.",
        ],
    )

    # 7. Model choices defended
    add_content_slide(
        prs,
        "Model choices — defended",
        [
            ("Red Team primary: huihui-ai/Llama-3.3-70B-Instruct-abliterated", [
                "Abliterated = refusal direction surgically removed from weights, no retraining. Reproducible technique, open weights on HF.",
                "Hosted on RunPod serverless GPU (A100-40GB, 4-bit quant, OpenAI-compatible API). Scales to zero when idle.",
                "Frontier-provider hosts (Together, Anthropic, OpenAI) will not host abliterated weights. RunPod / Modal / Hyperbolic will.",
            ]),
            ("Red Team escalation: DeepSeek-R1 via API", [
                "Materially lower refusal rate than Western frontier models on offensive prompts; ~10× cheaper than GPT-4 / Claude Sonnet.",
            ]),
            ("Judge: Claude Haiku 4.5", ["Cheap, fast, already validated as judge in Week-2 eval suite (24/25 baseline accuracy)."]),
            ("Orchestrator + Documentation: Claude Sonnet 4.6", ["Strategic reasoning / structured prose, no attack generation, no refusal problem."]),
        ],
    )

    # 8. Why PyRIT
    add_content_slide(
        prs,
        "Framework choice: PyRIT for attack orchestration",
        [
            ("Multi-turn mutation orchestrators built-in.", [
                "TreeOfAttacksWithPruningOrchestrator (TAP): attacker LLM generates → tests → refines based on target responses.",
                "CrescendoOrchestrator: gradual multi-turn escalation. ~1000 lines we'd otherwise hand-write and tune.",
            ]),
            ("Battle-tested.", [
                "Microsoft's AI Red Team uses it against Copilot, Phi-3, 100+ internal operations. Published paper, MIT-licensed, ~3.8K stars.",
            ]),
            ("Pluggable target/attacker models.", [
                "Swap RunPod-hosted abliterated Llama → DeepSeek-R1 → anything else by config. Same orchestrator, any OpenAI-compatible endpoint.",
            ]),
            ("Built-in converters (prompt obfuscation).", [
                "base64, leetspeak, ROT13, Unicode escape, character substitution — N × M coverage of every seed for free.",
            ]),
            ("Memory + scoring scaffolding.", [
                "Conversation memory, replay, and a scorer interface where our own Judge plugs in.",
            ]),
            ("Considered and rejected: Promptfoo (static eval, no multi-turn mutation), Garak (one-shot probing), rolling our own (~2 days reinventing TAP).", []),
        ],
        footer="What we don't use from PyRIT: high-level agent coordination (LangGraph), Judge wrapping (own rubrics), or its CLI patterns (we wrap it as a service).",
    )

    # 9. Orchestration strategy
    add_content_slide(
        prs,
        "Orchestration strategy — what gets tested next",
        [
            ("Coverage matrix: 6 categories × ~30 subcategories, each cell tracks cases-run, recent pass-rate, severity weight.", []),
            ("Priority score per cell = severity × (1 − coverage_ratio) + 0.4 × recent_failure_rate + 0.3 × time_since_last.", []),
            ("Campaign halts when 30-attack rolling success rate < 2% AND > $5 spent.", []),
            ("Hard caps enforced in harness, not trusted to agent prompts:", [
                "Per-campaign USD cap. Per-day global USD cap. Per-target QPS token-bucket.",
            ]),
            ("Regression runs auto-trigger on target deploys (CircleCI pre-`hold-qa` job calls /regression-runs).", []),
        ],
    )

    # 9. Regression & validation
    add_content_slide(
        prs,
        "Regression harness — deterministic where it counts",
        [
            ("Two assertion types:", [
                "DETERMINISTIC (preferred): regex on response, structural JSON check, tool-call audit, PHI byte-equality. Cheap, fast, no drift.",
                "JUDGED (only when semantic): re-runs the original Judge prompt. Tracked separately because of drift risk.",
            ]),
            ("Anti-Goodhart controls:", [
                "Every case carries a negative control — a benign prompt that must keep passing. Invalidates 'pass' if it breaks.",
                "Promotion stores the original rubric version. Drift > 2 minor versions → re-run with original rubric + report both.",
            ]),
            ("Spec gotcha avoided: a test that 'passes because model behavior changed' is worse than no test.", []),
        ],
    )

    # 10. CI/CD integration
    add_content_slide(
        prs,
        "CI/CD integration — gating qa/prod promotions",
        [
            ("Pre-`hold-qa` CircleCI job calls deployed adversary service's POST /regression-runs.", []),
            ("Dev auto-deploys ungated (fast inner loop preserved). qa and prod promotions are gated.", []),
            ("Fail if: any new HIGH-severity regression, OR overall pass rate drops >5%, OR cost-per-cycle on target rises >10%.", []),
            ("Bounded to ~3-5 min — promotion-gate suite is the deterministic subset; full LLM campaign runs async.", []),
            ("Emergency bypass: `[adversarial-bypass]` commit-msg convention + audit-trail-logged justification.", []),
        ],
        footer="Implementation Option A: CircleCI calls the adversary service via API. Suite is pinned to a tag, not main.",
    )

    # 11. UI for human operators (ad-hoc + dashboards)
    add_content_slide(
        prs,
        "Human-facing UI — Next.js dashboard",
        [
            ("Two operator paths into the platform: CI auto-runs (previous slide) AND a human UI.", []),
            ("Stack: Next.js 15 (App Router) + TypeScript + Tailwind + shadcn/ui + TanStack Query + recharts. SSE for live verdict streaming.", []),
            ("Pages — each maps to a user role from USERS.md:", [
                "/ Dashboard — open findings, last-24h runs, coverage heatmap, current spend (Security Engineer)",
                "/run — Ad hoc campaign: target + categories + budget → kick off → live verdict stream (Security Engineer)",
                "/coverage — 6×30 matrix of category × subcategory with case counts, pass-rate, time-since-last (Security Engineer, CISO)",
                "/findings + /findings/VULN-NNNN — searchable list, drafts queue with approve/reject gate, single-finding repro (All)",
                "/orchestrator — pause/resume, bump priorities, change budget caps (Security Engineer)",
                "/runs/<campaign_id> — every attack, every verdict, jump to Langfuse trace (All)",
                "/dashboard/exec — resilience trend lines + audit export (CISO)",
            ]),
            ("Standalone Next.js app, not iframe-embedded — this is a separate application from the W2 OpenEMR frontend.", []),
            ("Deployed as `adversary-ui` Railway service. SSE endpoint backed by FastAPI on `adversary-agent`.", []),
        ],
        footer="The CI path makes promotions safer. The UI path makes findings actionable — same data substrate (Postgres + Langfuse), different consumers.",
    )

    # 12. Trust & safety for the platform itself
    add_content_slide(
        prs,
        "Trust & safety for the platform itself",
        [
            ("Target-host allowlist (harness/executor.py) — defense of ARCHITECTURE.md §13:", [
                "Authorized: copilot-agent-production-41de.up.railway.app + qa + dev + localhost. Window: 2026-05-11 → 2026-05-22.",
                "Any HTTP call outside the allowlist hard-errors before bytes leave the platform.",
            ]),
            ("Human gates kept exactly where they belong:", [
                "Critical/high finding publication · rubric changes · allowlist edits · prod-target campaigns above $ cap · `[adversarial-bypass]` force-promote.",
            ]),
            ("Autonomy preserved exactly where it should be:", [
                "Campaign scheduling · attack mutation · individual judgments · medium/low finding drafts · nightly regressions · coverage priorities.",
            ]),
            ("No autonomous remediation — spec is explicit, and an agent that can push fixes can introduce vulnerabilities.", []),
            ("Immutable audit log on every campaign, attack, verdict, publish.", []),
        ],
    )

    # 12. Observability
    add_content_slide(
        prs,
        "Observability — Langfuse, tagged by agent_role",
        [
            "Self-hosted Langfuse (reuse W2 deployment) — no offensive-security trace leaves the platform.",
            ("Every call carries: agent_role, campaign_id, attack_category, attack_subcategory, target_sha, verdict.", []),
            ("Six required questions, all answerable from the dashboard:", [
                "Which categories have been tested, how many cases each?",
                "Pass/fail rate by category and target version?",
                "Is the target getting more or less resilient over time?",
                "Open vs in-progress vs resolved findings?",
                "Run cost, scaling rate?",
                "What did each agent do, in what order?",
            ]),
            "Per-call USD computed from (model, prompt_tokens, completion_tokens) stored on the span.",
        ],
    )

    # 13. Cost at scale
    add_content_slide(
        prs,
        "Cost at scale (~5% pass rate, 2:1 mutation:seed ratio)",
        [
            ("100 runs → ~$0.70 — no bottleneck.", []),
            ("1K runs → ~$7 — bottleneck: local Ollama / RunPod throughput.", []),
            ("10K runs → ~$70 — bottleneck: Judge concurrency / Anthropic rate limits.", []),
            ("100K runs → ~$700 base — required architectural changes:", [
                "Deterministic pre-filter knocks out ~40% of clearly-failed attacks before Judge sees them.",
                "Batched Judge calls + sharded Judge instances (one per category).",
                "Sharded RunPod / always-warm fleet for the 70B abliterated model.",
            ]),
            ("Cost-per-cycle is dominated by Judge, not Red Team — only because we picked uncensored local for generation.", []),
            ("Full breakdown in AI_COST_ANALYSIS.md.", []),
        ],
    )

    # 14. Sprint plan
    add_content_slide(
        prs,
        "Sprint plan (W3 deadlines)",
        [
            ("Today (Mon 05-11): architecture defense + THREAT_MODEL.md + USERS.md + initial seeds + Red Team prototype hitting live dev target.", []),
            ("Tue 05-12 (MVP): all four agents end-to-end on ≥3 categories, ≥3 vulnerability reports, platform deployed, observability MVP, cost analysis v1.", []),
            ("Wed 05-13: Orchestrator + coverage matrix + budget caps + regression harness wired to CircleCI pre-`hold-qa`.", []),
            ("Thu 05-14: Documentation Agent + full regression integration + UI dashboard polish + bug bash.", []),
            ("Fri 05-15 noon (FINAL): demo video, cost projections at 100/1K/10K/100K, README pass, social post.", []),
            ("Authorization window extends to 2026-05-22 for post-final exploration / replay.", []),
        ],
    )

    # 15. What we're prepared to defend
    add_content_slide(
        prs,
        "What I'm prepared to defend",
        [
            ("Choosing an uncensored local model for Red Team generation over Claude/GPT/Gemini.", []),
            ("Putting Judge on a frontier model and Red Team on a non-frontier model — opposite alignment requirements.", []),
            ("Postgres-backed inter-agent messaging instead of a broker — durability > latency at our volume.", []),
            ("Two frameworks (LangGraph + PyRIT) — LangGraph already in W2, PyRIT mutation strategies too valuable to skip.", []),
            ("Deterministic regression assertions wherever possible — judge drift would invalidate the entire history.", []),
            ("Human gate at high/critical findings, force-promote with `[adversarial-bypass]` audit trail.", []),
            ("Self-hosted Langfuse — offensive-security traces don't leave the platform.", []),
        ],
        footer="The deliverable that matters: one a hospital CISO would trust to continuously test systems their physicians depend on.",
    )

    out = "ARCHITECTURE_DEFENSE.pptx"
    prs.save(out)
    print(f"wrote {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
